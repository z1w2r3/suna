#!/usr/bin/env python3
"""
HTML Presentation to PPTX Converter - Perfect 1:1 Approach

This script provides PERFECT 1:1 conversion by:
- Capturing the ENTIRE slide as a pixel-perfect background image (including all icons, gradients, decorations)
- Making text transparent for the background capture
- Extracting text elements separately for editable PowerPoint text boxes
- Overlaying editable text on the perfect background

Usage:
    python html_to_pptx_perfect.py [presentation_directory] [output_pptx_path]

Example:
    python html_to_pptx_perfect.py . perfect_presentation.pptx
"""

import json
import os
import sys
import re
import asyncio
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import subprocess
from dataclasses import dataclass
import base64
import io

try:
    from playwright.async_api import async_playwright
except ImportError:
    print("Error: Playwright is not installed. Please install it with:")
    print("pip install playwright")
    print("playwright install chromium")
    sys.exit(1)

try:
    from bs4 import BeautifulSoup, Tag
except ImportError:
    print("Error: BeautifulSoup is not installed. Please install it with:")
    print("pip install beautifulsoup4")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print("Error: python-pptx is not installed or has missing components. Please install it with:")
    print("pip install python-pptx")
    print(f"Import error: {e}")
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("Error: Pillow is not installed. Please install it with:")
    print("pip install Pillow")
    sys.exit(1)


@dataclass
class TextElement:
    """Text element information for editable text boxes"""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_family: str
    font_size: float
    font_weight: str
    color: str
    text_align: str
    line_height: float
    tag: str


class CSSParser:
    """Parse CSS styles and convert values to appropriate units"""
    
    @staticmethod
    def parse_color(color_str: str) -> Tuple[int, int, int]:
        """Parse CSS color string to RGB tuple"""
        if not color_str:
            return (0, 0, 0)
        
        color_str = color_str.strip().lower()
        
        # Handle hex colors
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            try:
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            except ValueError:
                return (0, 0, 0)
        
        # Handle rgb() colors
        rgb_match = re.match(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
        if rgb_match:
            return tuple(int(x) for x in rgb_match.groups())
        
        # Handle rgba() colors (ignore alpha for now)
        rgba_match = re.match(r'rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_str)
        if rgba_match:
            return tuple(int(x) for x in rgba_match.groups())
        
        # Named colors
        named_colors = {
            'black': (0, 0, 0), 'white': (255, 255, 255), 'red': (255, 0, 0),
            'green': (0, 128, 0), 'blue': (0, 0, 255), 'yellow': (255, 255, 0),
            'cyan': (0, 255, 255), 'magenta': (255, 0, 255), 'gray': (128, 128, 128),
            'grey': (128, 128, 128), 'orange': (255, 165, 0), 'purple': (128, 0, 128)
        }
        
        return named_colors.get(color_str, (0, 0, 0))
    
    @staticmethod
    def parse_font_weight(weight_str: str) -> bool:
        """Parse font weight to bold boolean"""
        if not weight_str:
            return False
        
        weight_str = weight_str.strip().lower()
        bold_weights = ['bold', 'bolder', '700', '800', '900']
        
        return weight_str in bold_weights or (weight_str.isdigit() and int(weight_str) >= 700)


class PerfectHTMLToPPTXConverter:
    def __init__(self, presentation_dir: str, output_path: str = None):
        """
        Initialize the perfect converter.
        
        Args:
            presentation_dir: Directory containing metadata.json and HTML slides
            output_path: Output PPTX file path (optional)
        """
        self.presentation_dir = Path(presentation_dir).resolve()
        self.metadata_path = self.presentation_dir / "metadata.json"
        self.output_path = output_path
        self.metadata = None
        self.slides_info = []
        
        # Validate inputs
        if not self.presentation_dir.exists():
            raise FileNotFoundError(f"Presentation directory not found: {self.presentation_dir}")
        
        if not self.metadata_path.exists():
            raise FileNotFoundError(f"metadata.json not found in: {self.presentation_dir}")
    
    def load_metadata(self) -> Dict:
        """Load and parse metadata.json"""
        try:
            with open(self.metadata_path, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
            
            # Extract slide information and sort by slide number
            slides = self.metadata.get('slides', {})
            self.slides_info = []
            
            for slide_num, slide_data in slides.items():
                filename = slide_data.get('filename')
                title = slide_data.get('title', f'Slide {slide_num}')
                
                if filename:
                    html_path = self.presentation_dir / filename
                    if html_path.exists():
                        self.slides_info.append({
                            'number': int(slide_num),
                            'title': title,
                            'filename': filename,
                            'path': html_path
                        })
                    else:
                        print(f"Warning: HTML file not found: {html_path}")
            
            # Sort slides by number
            self.slides_info.sort(key=lambda x: x['number'])
            
            if not self.slides_info:
                raise ValueError("No valid slides found in metadata.json")
            
            # Set default output path if not provided
            if not self.output_path:
                presentation_name = self.metadata.get('presentation_name', 'presentation')
                self.output_path = self.presentation_dir / f"{presentation_name}_perfect.pptx"
            else:
                self.output_path = Path(self.output_path).resolve()
            
            print(f"Loaded {len(self.slides_info)} slides from metadata")
            return self.metadata
            
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in metadata.json: {e}")
        except Exception as e:
            raise ValueError(f"Error loading metadata: {e}")
    
    async def capture_perfect_background(self, browser, html_path: Path, temp_dir: Path) -> Path:
        """
        Capture the entire slide as a perfect background image with text made transparent.
        
        Args:
            browser: Playwright browser instance
            html_path: Path to HTML file
            temp_dir: Temporary directory for images
            
        Returns:
            Path to the perfect background image
        """
        page = await browser.new_page()
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1 for exact measurements
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Navigate to HTML file
            file_url = f"file://{html_path.absolute()}"
            await page.goto(file_url, wait_until="networkidle", timeout=30000)
            
            # Wait for fonts and content to load
            await page.wait_for_timeout(5000)
            
            # Make ALL text transparent while preserving layout and everything else
            await page.evaluate(r"""
                () => {
                    // Function to make text transparent while keeping all visual elements
                    function makeTextTransparent(element) {
                        if (element.nodeType === Node.TEXT_NODE) {
                            // Don't remove text nodes, just make them invisible
                            return;
                        } else if (element.nodeType === Node.ELEMENT_NODE) {
                            const computed = window.getComputedStyle(element);
                            
                            // If this element contains text, make the text transparent
                            // but preserve all other styling (backgrounds, borders, etc.)
                            const hasText = element.textContent && element.textContent.trim();
                            if (hasText) {
                                // Store original color for later if needed
                                const originalColor = computed.color;
                                element.setAttribute('data-original-color', originalColor);
                                
                                // Make text transparent but keep everything else
                                element.style.color = 'transparent';
                                element.style.textShadow = 'none';
                                element.style.webkitTextStroke = 'none';
                            }
                            
                            // Process children
                            Array.from(element.children).forEach(makeTextTransparent);
                        }
                    }
                    
                    // Apply to entire document
                    makeTextTransparent(document.body);
                    
                    console.log('Made all text transparent while preserving visual elements');
                }
            """)
            
            # Wait for changes to apply
            await page.wait_for_timeout(2000)
            
            # Take perfect screenshot
            background_path = temp_dir / f"perfect_background_{html_path.stem}.png"
            await page.screenshot(
                path=str(background_path),
                full_page=False,
                clip={"x": 0, "y": 0, "width": 1920, "height": 1080}
            )
            
            print(f"    ✓ Captured perfect background: {background_path.name}")
            return background_path
            
        except Exception as e:
            raise RuntimeError(f"Error capturing perfect background: {e}")
        finally:
            await page.close()
    
    async def extract_text_elements(self, browser, html_path: Path) -> List[TextElement]:
        """
        Extract all text elements with precise positioning for editable text boxes.
        
        Args:
            browser: Playwright browser instance
            html_path: Path to HTML file
            
        Returns:
            List of TextElement objects
        """
        page = await browser.new_page()
        text_elements = []
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Navigate to HTML file
            file_url = f"file://{html_path.absolute()}"
            await page.goto(file_url, wait_until="networkidle", timeout=30000)
            
            # Wait for fonts and content to load
            await page.wait_for_timeout(5000)
            
            # Extract all text elements with precise positioning
            text_data = await page.evaluate(r"""
                () => {
                    function extractTextFromElement(element) {
                        if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
                        
                        const computedStyle = window.getComputedStyle(element);
                        const rect = element.getBoundingClientRect();
                        
                        // Skip hidden elements
                        if (computedStyle.display === 'none' || computedStyle.visibility === 'hidden') return [];
                        if (rect.width === 0 || rect.height === 0) return [];
                        
                        const results = [];
                        
                        // Get direct text content (not from children)
                        let directText = '';
                        for (let node of element.childNodes) {
                            if (node.nodeType === Node.TEXT_NODE) {
                                directText += node.textContent;
                            }
                        }
                        directText = directText.trim();
                        
                        // If this element has direct text content, extract it
                        if (directText && directText.length > 0) {
                            const fontSizeMatch = computedStyle.fontSize.match(/([0-9.]+)px/);
                            const actualFontSize = fontSizeMatch ? parseFloat(fontSizeMatch[1]) : 16;
                            
                            results.push({
                                text: directText,
                                x: Math.round(rect.left * 100) / 100,
                                y: Math.round(rect.top * 100) / 100,
                                width: Math.round(rect.width * 100) / 100,
                                height: Math.round(rect.height * 100) / 100,
                                fontFamily: computedStyle.fontFamily,
                                actualFontSizePx: actualFontSize,
                                fontWeight: computedStyle.fontWeight,
                                color: computedStyle.color,
                                textAlign: computedStyle.textAlign,
                                lineHeight: computedStyle.lineHeight,
                                tag: element.tagName.toLowerCase()
                            });
                        }
                        
                        // Process children for nested text
                        Array.from(element.children).forEach(child => {
                            results.push(...extractTextFromElement(child));
                        });
                        
                        return results;
                    }
                    
                    // Start extraction from body
                    const allTextElements = extractTextFromElement(document.body);
                    
                    // Sort by position (top to bottom, left to right)
                    allTextElements.sort((a, b) => {
                        if (Math.abs(a.y - b.y) < 5) {  // Same line
                            return a.x - b.x;
                        }
                        return a.y - b.y;
                    });
                    
                    return allTextElements;
                }
            """)
            
            # Convert to TextElement objects
            for data in text_data:
                if data and data['text']:
                    # Parse font family
                    font_family = data['fontFamily']
                    if font_family:
                        font_family = font_family.split(',')[0].strip().strip('"\'')
                        font_family_map = {
                            'roboto': 'Roboto', 'arial': 'Arial', 'helvetica': 'Helvetica',
                            'sans-serif': 'Arial', 'serif': 'Times New Roman', 'monospace': 'Courier New'
                        }
                        font_family = font_family_map.get(font_family.lower(), font_family)
                    else:
                        font_family = 'Arial'
                    
                    # Parse line height
                    line_height = 1.2
                    if data['lineHeight'] and data['lineHeight'] != 'normal':
                        if data['lineHeight'].endswith('px'):
                            px_value = float(data['lineHeight'][:-2])
                            line_height = px_value / data['actualFontSizePx']
                        else:
                            try:
                                line_height = float(data['lineHeight'])
                            except:
                                line_height = 1.2
                    
                    text_element = TextElement(
                        text=data['text'],
                        x=data['x'],
                        y=data['y'],
                        width=data['width'],
                        height=data['height'],
                        font_family=font_family,
                        font_size=data['actualFontSizePx'] * 0.75,  # Convert px to points
                        font_weight=data['fontWeight'],
                        color=data['color'],
                        text_align=data['textAlign'],
                        line_height=line_height,
                        tag=data['tag']
                    )
                    
                    text_elements.append(text_element)
            
            print(f"    ✓ Extracted {len(text_elements)} text elements")
            return text_elements
            
        except Exception as e:
            raise RuntimeError(f"Error extracting text elements: {e}")
        finally:
            await page.close()
    
    def create_text_box(self, slide, text_element: TextElement) -> None:
        """
        Create an editable text box in PowerPoint with exact positioning.
        
        Args:
            slide: PowerPoint slide object
            text_element: TextElement with positioning and styling
        """
        # Convert pixel coordinates to inches
        left = Inches(text_element.x / 96.0)
        top = Inches(text_element.y / 96.0)
        width = Inches(max(text_element.width, 10) / 96.0)
        height = Inches(max(text_element.height, 10) / 96.0)
        
        print(f"    📝 Text: '{text_element.text[:40]}...' at ({text_element.x:.1f}, {text_element.y:.1f})")
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Set text frame properties for exact positioning
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.word_wrap = True
        text_frame.auto_size = None
        
        # Add paragraph
        p = text_frame.paragraphs[0]
        p.text = text_element.text
        
        # Set text alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'centre': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY, 'start': PP_ALIGN.LEFT,
            'end': PP_ALIGN.RIGHT
        }
        p.alignment = alignment_map.get(text_element.text_align.lower(), PP_ALIGN.LEFT)
        
        # Set spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if hasattr(p, 'line_spacing'):
            p.line_spacing = text_element.line_height
        
        # Set font properties
        font = p.font
        font.name = text_element.font_family
        font.size = Pt(max(text_element.font_size, 8))
        font.bold = CSSParser.parse_font_weight(text_element.font_weight)
        
        # Set font color
        try:
            r, g, b = CSSParser.parse_color(text_element.color)
            font.color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"    Warning: Could not parse color '{text_element.color}': {e}")
            font.color.rgb = RGBColor(0, 0, 0)
        
        # Make textbox transparent (no background, no border)
        textbox.fill.background()
        textbox.line.fill.background()
    
    async def convert_slide_perfect(self, browser, slide_info: Dict, presentation, temp_dir: Path) -> None:
        """
        Convert a single HTML slide using perfect 1:1 approach.
        
        Args:
            browser: Playwright browser instance
            slide_info: Slide information dictionary
            presentation: PowerPoint presentation object
            temp_dir: Temporary directory for images
        """
        html_path = slide_info['path']
        slide_num = slide_info['number']
        
        print(f"Converting slide {slide_num}: {slide_info['title']} (Perfect 1:1 Mode)")
        
        # Add blank slide
        blank_slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(blank_slide_layout)
        
        # Step 1: Capture perfect background (everything except text)
        print("  🎨 Capturing PERFECT background with all visual elements...")
        background_image_path = await self.capture_perfect_background(browser, html_path, temp_dir)
        
        # Step 2: Add perfect background to slide
        if background_image_path and background_image_path.exists():
            left = Inches(0)
            top = Inches(0)
            width = Inches(20)  # 1920px at 96 DPI
            height = Inches(11.25)  # 1080px at 96 DPI
            
            picture = slide.shapes.add_picture(str(background_image_path), left, top, width, height)
            print(f"    ✅ Perfect background added (1920x1080)")
        
        # Step 3: Extract and add editable text elements
        print("  📝 Extracting editable text elements...")
        text_elements = await self.extract_text_elements(browser, html_path)
        
        # Step 4: Create editable text boxes on top of perfect background
        print("  ✍️  Adding editable text overlays...")
        for text_element in text_elements:
            self.create_text_box(slide, text_element)
        
        print(f"  🎉 Slide {slide_num}: PERFECT background + {len(text_elements)} editable text elements")
    
    async def convert_to_pptx_perfect(self) -> None:
        """Main perfect conversion method"""
        print("🎯 Starting PERFECT 1:1 HTML to PPTX conversion...")
        print("📋 Method: Perfect background rasterization + Editable text overlay")
        print("=" * 80)
        
        # Load metadata
        self.load_metadata()
        
        # Create new PowerPoint presentation
        presentation = Presentation()
        
        # Set slide dimensions to 1920x1080 (16:9)
        presentation.slide_width = Inches(20)  # 1920px at 96 DPI
        presentation.slide_height = Inches(11.25)  # 1080px at 96 DPI
        
        # Remove default slide
        if len(presentation.slides) > 0:
            xml_slides = presentation.slides._sldIdLst
            xml_slides.remove(xml_slides[0])
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Launch browser for processing
            async with async_playwright() as p:
                print("🌐 Launching browser for perfect processing...")
                browser = await p.chromium.launch(
                    headless=True,
                    args=[
                        '--no-sandbox',
                        '--disable-setuid-sandbox',
                        '--disable-dev-shm-usage',
                        '--disable-gpu',
                        '--force-device-scale-factor=1',
                        '--disable-background-timer-throttling',
                        '--disable-backgrounding-occluded-windows',
                        '--disable-renderer-backgrounding',
                        '--disable-features=VizDisplayCompositor'
                    ]
                )
                
                try:
                    # Process each slide
                    for slide_info in self.slides_info:
                        await self.convert_slide_perfect(browser, slide_info, presentation, temp_path)
                    
                finally:
                    await browser.close()
        
        # Save PowerPoint presentation
        presentation.save(str(self.output_path))
        print(f"\n🎉 PERFECT 1:1 PPTX created successfully: {self.output_path}")
        print(f"📊 Total slides: {len(presentation.slides)}")
        print(f"✨ Perfect visual fidelity + Fully editable text!")


def check_dependencies():
    """Check if required dependencies are available"""
    missing_deps = []
    
    try:
        import playwright
    except ImportError:
        missing_deps.append("playwright (pip install playwright)")
    
    try:
        from pptx import Presentation
    except ImportError:
        missing_deps.append("python-pptx (pip install python-pptx)")
    
    try:
        from PIL import Image
    except ImportError:
        missing_deps.append("Pillow (pip install Pillow)")
    
    if missing_deps:
        print("❌ Missing dependencies:")
        for dep in missing_deps:
            print(f"  - {dep}")
        print("\nPlease install missing dependencies and try again.")
        return False
    
    return True


def main():
    """Main CLI entry point"""
    print("🎯 HTML Presentation to PPTX Converter - PERFECT 1:1 MODE")
    print("=" * 80)
    print("🎨 Perfect background capture + Editable text overlay")
    print()
    
    # Check dependencies
    if not check_dependencies():
        sys.exit(1)
    
    # Parse command line arguments
    if len(sys.argv) < 2:
        presentation_dir = "."
        output_path = None
    elif len(sys.argv) == 2:
        presentation_dir = sys.argv[1]
        output_path = None
    elif len(sys.argv) == 3:
        presentation_dir = sys.argv[1]
        output_path = sys.argv[2]
    else:
        print("Usage: python html_to_pptx_perfect.py [presentation_directory] [output_pptx_path]")
        print("\nExamples:")
        print("  python html_to_pptx_perfect.py")
        print("  python html_to_pptx_perfect.py . perfect_presentation.pptx")
        print("  python html_to_pptx_perfect.py /path/to/slides output.pptx")
        sys.exit(1)
    
    try:
        # Create converter and run
        converter = PerfectHTMLToPPTXConverter(presentation_dir, output_path)
        asyncio.run(converter.convert_to_pptx_perfect())
        
    except KeyboardInterrupt:
        print("\n❌ Conversion cancelled by user")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
