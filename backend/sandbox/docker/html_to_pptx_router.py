#!/usr/bin/env python3
"""
FastAPI HTML Presentation to PPTX Converter Router

Provides PPTX conversion endpoints as a FastAPI router that can be included in other applications.
Uses the perfect 1:1 conversion approach with async processing.
"""

import json
import asyncio
from pathlib import Path
from typing import Dict, List, Optional
import tempfile
import shutil
from dataclasses import dataclass

from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel, Field

try:
    from playwright.async_api import async_playwright
except ImportError:
    raise ImportError("Playwright is not installed. Please install it with: pip install playwright")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    raise ImportError(f"python-pptx is not installed. Please install it with: pip install python-pptx. Error: {e}")

try:
    from bs4 import BeautifulSoup
except ImportError:
    raise ImportError("BeautifulSoup is not installed. Please install it with: pip install beautifulsoup4")


# Create router
router = APIRouter(prefix="/presentation", tags=["pptx-conversion"])

# Create output directory for generated PPTXs
output_dir = Path("generated_pptx")
output_dir.mkdir(exist_ok=True)


class ConvertRequest(BaseModel):
    presentation_path: str = Field(..., description="Path to the presentation folder containing metadata.json")
    download: bool = Field(False, description="If true, returns the PPTX file directly. If false, returns JSON with download URL.")


class ConvertResponse(BaseModel):
    success: bool
    message: str
    pptx_url: str
    filename: str
    total_slides: int


@dataclass
class TextElement:
    """Text element information for editable text boxes"""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_family: str
    font_size: float
    font_weight: str
    color: str
    text_align: str
    line_height: float
    tag: str


class CSSParser:
    """Parse CSS styles and convert values to appropriate units"""
    
    @staticmethod
    def parse_color(color_str: str) -> tuple[int, int, int]:
        """Parse CSS color string to RGB tuple"""
        if not color_str:
            return (0, 0, 0)
        
        color_str = color_str.strip().lower()
        
        # Handle hex colors
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            try:
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            except ValueError:
                return (0, 0, 0)
        
        # Handle rgb() colors
        import re
        rgb_match = re.match(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
        if rgb_match:
            return tuple(int(x) for x in rgb_match.groups())
        
        # Handle rgba() colors (ignore alpha for now)
        rgba_match = re.match(r'rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_str)
        if rgba_match:
            return tuple(int(x) for x in rgba_match.groups())
        
        # Named colors
        named_colors = {
            'black': (0, 0, 0), 'white': (255, 255, 255), 'red': (255, 0, 0),
            'green': (0, 128, 0), 'blue': (0, 0, 255), 'yellow': (255, 255, 0),
            'cyan': (0, 255, 255), 'magenta': (255, 0, 255), 'gray': (128, 128, 128),
            'grey': (128, 128, 128), 'orange': (255, 165, 0), 'purple': (128, 0, 128)
        }
        
        return named_colors.get(color_str, (0, 0, 0))
    
    @staticmethod
    def parse_font_weight(weight_str: str) -> bool:
        """Parse font weight to bold boolean"""
        if not weight_str:
            return False
        
        weight_str = weight_str.strip().lower()
        bold_weights = ['bold', 'bolder', '700', '800', '900']
        
        return weight_str in bold_weights or (weight_str.isdigit() and int(weight_str) >= 700)


class PerfectHTMLToPPTXConverter:
    def __init__(self, presentation_dir: str):
        """Initialize the perfect converter."""
        self.presentation_dir = Path(presentation_dir).resolve()
        self.metadata_path = self.presentation_dir / "metadata.json"
        self.metadata = None
        self.slides_info = []
        
        # Validate inputs
        if not self.presentation_dir.exists():
            raise FileNotFoundError(f"Presentation directory not found: {self.presentation_dir}")
        
        if not self.metadata_path.exists():
            raise FileNotFoundError(f"metadata.json not found in: {self.presentation_dir}")
    
    def load_metadata(self) -> Dict:
        """Load and parse metadata.json"""
        try:
            with open(self.metadata_path, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
            
            # Extract slide information and sort by slide number
            slides = self.metadata.get('slides', {})
            self.slides_info = []
            
            for slide_num, slide_data in slides.items():
                filename = slide_data.get('filename')
                file_path = slide_data.get('file_path')
                title = slide_data.get('title', f'Slide {slide_num}')
                
                if file_path:
                    # Use the local presentation directory instead of hardcoded /workspace
                    html_path = Path(f"/workspace/{file_path}")
                    print(f"Using path: {html_path}")
                    
                    # Verify the path exists
                    if html_path.exists():
                        self.slides_info.append({
                            'number': int(slide_num),
                            'title': title,
                            'filename': filename,
                            'path': html_path
                        })
                        print(f"Added slide {slide_num}: {html_path}")
                    else:
                        print(f"Warning: HTML file does not exist: {html_path}")
            
            # Sort slides by number
            self.slides_info.sort(key=lambda x: x['number'])
            
            if not self.slides_info:
                raise ValueError("No valid slides found in metadata.json")
            
            return self.metadata
            
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in metadata.json: {e}")
        except Exception as e:
            raise ValueError(f"Error loading metadata: {e}")
    
    async def extract_visual_elements(self, browser, html_path: Path, temp_dir: Path) -> List[Dict]:
        """Extract all visual elements (non-text) as individual images with positioning."""
        page = await browser.new_page()
        visual_elements = []
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Navigate to HTML file
            file_url = f"file://{html_path.absolute()}"
            await page.goto(file_url, wait_until="networkidle", timeout=30000)
            
            # Wait for fonts and content to load
            await page.wait_for_timeout(5000)
            
            # Extract all visual elements with precise positioning
            visual_data = await page.evaluate(r"""
                () => {
                    function extractVisualElements(element, depth = 0) {
                        if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
                        
                        const computedStyle = window.getComputedStyle(element);
                        const rect = element.getBoundingClientRect();
                        
                        // Skip hidden elements
                        if (computedStyle.display === 'none' || computedStyle.visibility === 'hidden') return [];
                        if (rect.width === 0 || rect.height === 0) return [];
                        
                        const results = [];
                        
                        // Check if this element has visual content (not just text)
                        const hasVisualContent = (
                            computedStyle.backgroundImage !== 'none' ||
                            computedStyle.backgroundColor !== 'rgba(0, 0, 0, 0)' ||
                            computedStyle.borderStyle !== 'none' ||
                            computedStyle.borderImage !== 'none' ||
                            element.tagName === 'IMG' ||
                            element.tagName === 'SVG' ||
                            element.tagName === 'CANVAS' ||
                            computedStyle.boxShadow !== 'none' ||
                            computedStyle.borderRadius !== '0px' ||
                            computedStyle.gradient !== 'none'
                        );
                        
                        // Check if this element has meaningful text content
                        const hasText = element.textContent && element.textContent.trim().length > 0;
                        
                        // Skip very large elements that are likely the main background
                        const isLikelyBackground = rect.width > 1000 || rect.height > 800;
                        
                        // If element has visual content, is not just a text container, and is not the main background
                        if (hasVisualContent && !hasText && !isLikelyBackground) {
                            results.push({
                                type: 'visual',
                                x: Math.round(rect.left * 100) / 100,
                                y: Math.round(rect.top * 100) / 100,
                                width: Math.round(rect.width * 100) / 100,
                                height: Math.round(rect.height * 100) / 100,
                                tag: element.tagName.toLowerCase(),
                                className: element.className,
                                id: element.id,
                                depth: depth
                            });
                        }
                        
                        // Process children for nested elements
                        Array.from(element.children).forEach(child => {
                            results.push(...extractVisualElements(child, depth + 1));
                        });
                        
                        return results;
                    }
                    
                    // Start extraction from body
                    const allVisualElements = extractVisualElements(document.body);
                    
                    // Sort by depth (background elements first) then by position
                    allVisualElements.sort((a, b) => {
                        if (a.depth !== b.depth) return a.depth - b.depth;
                        if (Math.abs(a.y - b.y) < 5) return a.x - b.x;
                        return a.y - b.y;
                    });
                    
                    return allVisualElements;
                }
            """)
            
            # Capture each visual element as an individual image
            for i, data in enumerate(visual_data):
                if data and data['type'] == 'visual':
                    # Make text transparent for this element only
                    await page.evaluate(r"""
                        (elementData) => {
                            const elements = document.querySelectorAll('*');
                            for (let el of elements) {
                                const rect = el.getBoundingClientRect();
                                if (Math.abs(rect.left - elementData.x) < 1 && 
                                    Math.abs(rect.top - elementData.y) < 1 &&
                                    Math.abs(rect.width - elementData.width) < 1 &&
                                    Math.abs(rect.height - elementData.height) < 1) {
                                    
                                    // Make text transparent in this element
                                    if (el.textContent && el.textContent.trim()) {
                                        el.style.color = 'transparent';
                                        el.style.textShadow = 'none';
                                        el.style.webkitTextStroke = 'none';
                                    }
                                    break;
                                }
                            }
                        }
                    """, data)
                    
                    # Wait for changes to apply
                    await page.wait_for_timeout(500)
                    
                    # Capture this specific element
                    element_path = temp_dir / f"visual_element_{html_path.stem}_{i:03d}.png"
                    await page.screenshot(
                        path=str(element_path),
                        full_page=False,
                        clip={
                            "x": data['x'], 
                            "y": data['y'], 
                            "width": data['width'], 
                            "height": data['height']
                        }
                    )
                    
                    visual_element = {
                        'type': 'visual',
                        'x': data['x'],
                        'y': data['y'],
                        'width': data['width'],
                        'height': data['height'],
                        'tag': data['tag'],
                        'image_path': element_path,
                        'depth': data['depth']
                    }
                    
                    visual_elements.append(visual_element)
                    print(f"    ✓ Captured visual element {i}: {data['tag']} at ({data['x']:.1f}, {data['y']:.1f})")
            
            print(f"    ✓ Extracted {len(visual_elements)} visual elements")
            
            # Also capture the full slide background to ensure we don't miss any background colors
            print("    🎨 Capturing full slide background as fallback...")
            try:
                # Make all text AND visual elements transparent temporarily
                await page.evaluate(r"""
                    () => {
                        function makeAllContentTransparent(element) {
                            if (element.nodeType === Node.TEXT_NODE) return;
                            if (element.nodeType === Node.ELEMENT_NODE) {
                                // Make text transparent
                                if (element.textContent && element.textContent.trim()) {
                                    element.style.color = 'transparent';
                                    element.style.textShadow = 'none';
                                    element.style.webkitTextStroke = 'none';
                                }
                                
                                // Make visual elements transparent (icons, graphics, borders, etc.)
                                const computedStyle = window.getComputedStyle(element);
                                const hasVisualContent = (
                                    computedStyle.backgroundImage !== 'none' ||
                                    computedStyle.backgroundColor !== 'rgba(0, 0, 0, 0)' ||
                                    computedStyle.borderStyle !== 'none' ||
                                    computedStyle.borderImage !== 'none' ||
                                    element.tagName === 'IMG' ||
                                    element.tagName === 'SVG' ||
                                    element.tagName === 'CANVAS' ||
                                    computedStyle.boxShadow !== 'none' ||
                                    computedStyle.borderRadius !== '0px' ||
                                    computedStyle.gradient !== 'none'
                                );
                                
                                // Make visual elements transparent but keep their space
                                if (hasVisualContent) {
                                    element.style.opacity = '0';
                                    element.style.pointerEvents = 'none';
                                }
                                
                                // Process children
                                Array.from(element.children).forEach(makeAllContentTransparent);
                            }
                        }
                        makeAllContentTransparent(document.body);
                    }
                """)
                
                await page.wait_for_timeout(1000)
                
                # Capture full slide
                full_background_path = temp_dir / f"full_background_{html_path.stem}.png"
                await page.screenshot(
                    path=str(full_background_path),
                    full_page=False,
                    clip={"x": 0, "y": 0, "width": 1920, "height": 1080}
                )
                
                # Add full background as first visual element (will be placed behind everything)
                full_background_element = {
                    'type': 'background',
                    'x': 0,
                    'y': 0,
                    'width': 1920,
                    'height': 1080,
                    'tag': 'full_background',
                    'image_path': full_background_path,
                    'depth': -1  # Lowest depth to place behind everything
                }
                
                visual_elements.insert(0, full_background_element)
                print(f"    ✅ Added full background fallback")
                
            except Exception as e:
                print(f"    ⚠️  Could not capture full background: {e}")
            
            return visual_elements
            
        except Exception as e:
            raise RuntimeError(f"Error extracting visual elements: {e}")
        finally:
            await page.close()

    async def capture_clean_background(self, browser, html_path: Path, temp_dir: Path, visual_elements: List[Dict]) -> Path:
        """Capture the clean background with visual elements temporarily hidden."""
        page = await browser.new_page()
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1 for exact measurements
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Navigate to HTML file
            file_url = f"file://{html_path.absolute()}"
            await page.goto(file_url, wait_until="networkidle", timeout=30000)
            
            # Wait for fonts and content to load
            await page.wait_for_timeout(5000)
            
            # Make text transparent AND hide visual elements to get clean background
            await page.evaluate(r"""
                (visualElementsData) => {
                    // Function to make text transparent
                    function makeTextTransparent(element) {
                        if (element.nodeType === Node.TEXT_NODE) return;
                        if (element.nodeType === Node.ELEMENT_NODE) {
                            const hasText = element.textContent && element.textContent.trim();
                            if (hasText) {
                                element.style.color = 'transparent';
                                element.style.textShadow = 'none';
                                element.style.webkitTextStroke = 'none';
                            }
                            Array.from(element.children).forEach(makeTextTransparent);
                        }
                    }
                    
                    // Function to hide visual elements (but preserve main background)
                    function hideVisualElements() {
                        visualElementsData.forEach(visualData => {
                            // Skip very large elements that are likely the main background
                            if (visualData.width > 1000 || visualData.height > 800) {
                                console.log('Skipping large element (likely background):', visualData);
                                return;
                            }
                            
                            // Find elements that match the visual element criteria
                            const elements = document.querySelectorAll('*');
                            for (let el of elements) {
                                const rect = el.getBoundingClientRect();
                                if (Math.abs(rect.left - visualData.x) < 2 && 
                                    Math.abs(rect.top - visualData.y) < 2 &&
                                    Math.abs(rect.width - visualData.width) < 2 &&
                                    Math.abs(rect.height - visualData.height) < 2) {
                                    
                                    // Hide this visual element
                                    el.style.visibility = 'hidden';
                                    el.setAttribute('data-temporarily-hidden', 'true');
                                    break;
                                }
                            }
                        });
                    }
                    
                    // Apply both transformations
                    makeTextTransparent(document.body);
                    hideVisualElements();
                    
                    console.log('Made text transparent and hid visual elements for clean background');
                }
            """, visual_elements)
            
            # Wait for changes to apply
            await page.wait_for_timeout(2000)
            
            # Take clean background screenshot
            background_path = temp_dir / f"clean_background_{html_path.stem}.png"
            await page.screenshot(
                path=str(background_path),
                full_page=False,
                clip={"x": 0, "y": 0, "width": 1920, "height": 1080}
            )
            
            # Verify the captured image
            if background_path.exists():
                file_size = background_path.stat().st_size
                print(f"    ✓ Captured clean background: {background_path.name} ({file_size} bytes)")
                
                # Debug: check if image is mostly black
                try:
                    from PIL import Image
                    with Image.open(background_path) as img:
                        # Convert to RGB if needed
                        if img.mode != 'RGB':
                            img = img.convert('RGB')
                        
                        # Sample some pixels to check if image is mostly black
                        pixels = list(img.getdata())
                        black_pixels = sum(1 for p in pixels if p[0] < 50 and p[1] < 50 and p[2] < 50)
                        total_pixels = len(pixels)
                        black_percentage = (black_pixels / total_pixels) * 100
                        
                        print(f"    📊 Image analysis: {black_percentage:.1f}% black pixels")
                        
                        if black_percentage > 80:
                            print(f"    ⚠️  WARNING: Image appears to be mostly black!")
                        
                except Exception as e:
                    print(f"    ⚠️  Could not analyze image: {e}")
            else:
                print(f"    ❌ Background image was not created!")
            
            return background_path
            
        except Exception as e:
            raise RuntimeError(f"Error capturing clean background: {e}")
        finally:
            await page.close()
    
    async def extract_text_elements(self, browser, html_path: Path) -> List[TextElement]:
        """Extract all text elements with precise positioning for editable text boxes."""
        page = await browser.new_page()
        text_elements = []
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Navigate to HTML file
            file_url = f"file://{html_path.absolute()}"
            await page.goto(file_url, wait_until="networkidle", timeout=30000)
            
            # Wait for fonts and content to load
            await page.wait_for_timeout(5000)
            
            # Extract all text elements with precise positioning
            text_data = await page.evaluate(r"""
                () => {
                    function extractTextFromElement(element) {
                        if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
                        
                        const computedStyle = window.getComputedStyle(element);
                        const rect = element.getBoundingClientRect();
                        
                        // Skip hidden elements
                        if (computedStyle.display === 'none' || computedStyle.visibility === 'hidden') return [];
                        if (rect.width === 0 || rect.height === 0) return [];
                        
                        const results = [];
                        
                        // Get direct text content (not from children)
                        let directText = '';
                        for (let node of element.childNodes) {
                            if (node.nodeType === Node.TEXT_NODE) {
                                directText += node.textContent;
                            }
                        }
                        directText = directText.trim();
                        
                        // If this element has direct text content, extract it
                        if (directText && directText.length > 0) {
                            const fontSizeMatch = computedStyle.fontSize.match(/([0-9.]+)px/);
                            const actualFontSize = fontSizeMatch ? parseFloat(fontSizeMatch[1]) : 16;
                            
                            // Get the actual text node position, not the container
                            const textNodes = [];
                            for (let node of element.childNodes) {
                                if (node.nodeType === Node.TEXT_NODE && node.textContent.trim()) {
                                    const range = document.createRange();
                                    range.selectNodeContents(node);
                                    const textRect = range.getBoundingClientRect();
                                    textNodes.push({
                                        text: node.textContent.trim(),
                                        rect: textRect
                                    });
                                }
                            }
                            
                            // Use text node position if available, otherwise fall back to container
                            if (textNodes.length > 0) {
                                textNodes.forEach(textNode => {
                                    results.push({
                                        text: textNode.text,
                                        x: Math.round(textNode.rect.left * 100) / 100,
                                        y: Math.round(textNode.rect.top * 100) / 100,
                                        width: Math.round(textNode.rect.width * 100) / 100,
                                        height: Math.round(textNode.rect.height * 100) / 100,
                                        fontFamily: computedStyle.fontFamily,
                                        actualFontSizePx: actualFontSize,
                                        fontWeight: computedStyle.fontWeight,
                                        color: computedStyle.color,
                                        textAlign: computedStyle.textAlign,
                                        lineHeight: computedStyle.lineHeight,
                                        tag: element.tagName.toLowerCase()
                                    });
                                });
                            } else {
                                // Fallback to container position
                                results.push({
                                    text: directText,
                                    x: Math.round(rect.left * 100) / 100,
                                    y: Math.round(rect.top * 100) / 100,
                                    width: Math.round(rect.width * 100) / 100,
                                    height: Math.round(rect.height * 100) / 100,
                                    fontFamily: computedStyle.fontFamily,
                                    actualFontSizePx: actualFontSize,
                                    fontWeight: computedStyle.fontWeight,
                                    color: computedStyle.color,
                                    textAlign: computedStyle.textAlign,
                                    lineHeight: computedStyle.lineHeight,
                                    tag: element.tagName.toLowerCase()
                                });
                            }
                        }
                        
                        // Process children for nested text
                        Array.from(element.children).forEach(child => {
                            results.push(...extractTextFromElement(child));
                        });
                        
                        return results;
                    }
                    
                    // Start extraction from body
                    const allTextElements = extractTextFromElement(document.body);
                    
                    // Sort by position (top to bottom, left to right)
                    allTextElements.sort((a, b) => {
                        if (Math.abs(a.y - b.y) < 5) {  // Same line
                            return a.x - b.x;
                        }
                        return a.y - b.y;
                    });
                    
                    return allTextElements;
                }
            """)
            
            # Convert to TextElement objects
            for data in text_data:
                if data and data['text']:
                    # Parse font family
                    font_family = data['fontFamily']
                    if font_family:
                        font_family = font_family.split(',')[0].strip().strip('"\'')
                        font_family_map = {
                            'roboto': 'Roboto', 'arial': 'Arial', 'helvetica': 'Helvetica',
                            'sans-serif': 'Arial', 'serif': 'Times New Roman', 'monospace': 'Courier New'
                        }
                        font_family = font_family_map.get(font_family.lower(), font_family)
                    else:
                        font_family = 'Arial'
                    
                    # Parse line height
                    line_height = 1.2
                    if data['lineHeight'] and data['lineHeight'] != 'normal':
                        if data['lineHeight'].endswith('px'):
                            px_value = float(data['lineHeight'][:-2])
                            line_height = px_value / data['actualFontSizePx']
                        else:
                            try:
                                line_height = float(data['lineHeight'])
                            except:
                                line_height = 1.2
                    
                    text_element = TextElement(
                        text=data['text'],
                        x=data['x'],
                        y=data['y'],
                        width=data['width'],
                        height=data['height'],
                        font_family=font_family,
                        font_size=data['actualFontSizePx'] * 0.75,  # Convert px to points
                        font_weight=data['fontWeight'],
                        color=data['color'],
                        text_align=data['textAlign'],
                        line_height=line_height,
                        tag=data['tag']
                    )
                    
                    text_elements.append(text_element)
            
            print(f"    ✓ Extracted {len(text_elements)} text elements")
            return text_elements
            
        except Exception as e:
            raise RuntimeError(f"Error extracting text elements: {e}")
        finally:
            await page.close()
    
    def create_text_box(self, slide, text_element: TextElement) -> None:
        """Create an editable text box in PowerPoint with exact positioning."""
        # Convert pixel coordinates to inches
        left = Inches(text_element.x / 96.0)
        top = Inches(text_element.y / 96.0)
        
        # Calculate proper width - ensure it's wide enough for the text content
        # Add some padding to prevent text wrapping
        text_width = text_element.width
        if text_width < 50:  # If width is too small, estimate based on text length
            # Rough estimate: each character needs about 8-12px depending on font
            estimated_width = len(text_element.text) * 10  # 10px per character
            text_width = max(text_width, estimated_width)
        
        # Add padding to prevent text from being cut off
        padding = 8  # 8px padding on each side
        final_width = text_width + (padding * 2)
        
        width = Inches(final_width / 96.0)
        height = Inches(max(text_element.height, 10) / 96.0)
        
        print(f"    📝 Text: '{text_element.text[:40]}...' at ({text_element.x:.1f}, {text_element.y:.1f})")
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Set text frame properties for exact positioning
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.word_wrap = True
        text_frame.auto_size = None
        
        # Add paragraph
        p = text_frame.paragraphs[0]
        p.text = text_element.text
        
        # Set text alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'centre': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY, 'start': PP_ALIGN.LEFT,
            'end': PP_ALIGN.RIGHT
        }
        p.alignment = alignment_map.get(text_element.text_align.lower(), PP_ALIGN.LEFT)
        
        # Set spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if hasattr(p, 'line_spacing'):
            p.line_spacing = text_element.line_height
        
        # Set font properties
        font = p.font
        font.name = text_element.font_family
        font.size = Pt(max(text_element.font_size, 8))
        font.bold = CSSParser.parse_font_weight(text_element.font_weight)
        
        # Set font color
        try:
            r, g, b = CSSParser.parse_color(text_element.color)
            font.color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"    Warning: Could not parse color '{text_element.color}': {e}")
            font.color.rgb = RGBColor(0, 0, 0)
        
        # Make textbox transparent (no background, no border)
        textbox.fill.background()
        textbox.line.fill.background()
    
    def add_visual_element_to_slide(self, slide, visual_element: Dict) -> None:
        """Add a visual element as an image to the PowerPoint slide with exact positioning."""
        if visual_element['image_path'].exists():
            # Convert pixel coordinates to inches
            left = Inches(visual_element['x'] / 96.0)
            top = Inches(visual_element['y'] / 96.0)
            width = Inches(visual_element['width'] / 96.0)
            height = Inches(visual_element['height'] / 96.0)
            
            # Add the image to the slide
            picture = slide.shapes.add_picture(str(visual_element['image_path']), left, top, width, height)
            
            # Special handling for background elements
            if visual_element['tag'] == 'clean_background':
                # Ensure background is sent to back
                picture.z_order = 0
                print(f"    🎨 Added background: {visual_element['image_path'].name} (z-order: {picture.z_order})")
            else:
                print(f"    ✅ Added visual element: {visual_element['tag']} at ({visual_element['x']:.1f}, {visual_element['y']:.1f})")
        else:
            print(f"    ⚠️  Visual element image not found: {visual_element['image_path']}")
            # Debug: check what files exist in temp_dir
            temp_dir = visual_element['image_path'].parent
            if temp_dir.exists():
                print(f"    📁 Files in temp directory: {list(temp_dir.glob('*.png'))}")
    
    async def build_slide_from_analysis(self, presentation, slide_analysis: Dict, temp_dir: Path) -> None:
        """Build a PowerPoint slide from pre-analyzed data (fast, no browser operations)."""
        slide_info = slide_analysis['slide_info']
        visual_elements = slide_analysis['visual_elements']
        background_path = slide_analysis['background_path']
        text_elements = slide_analysis['text_elements']
        
        slide_num = slide_info['number']
        print(f"  🏗️  Building slide {slide_num}: {slide_info['title']}")
        
        # Add blank slide
        blank_slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(blank_slide_layout)
        
        # Step 1: Add the clean background as the base layer
        print(f"    🖼️  Adding background...")
        background_element = {
            'type': 'background',
            'x': 0,
            'y': 0,
            'width': 1920,
            'height': 1080,
            'tag': 'clean_background',
            'image_path': background_path,
            'depth': -1
        }
        self.add_visual_element_to_slide(slide, background_element)
        
        # Step 2: Add individual visual elements on top of background
        print(f"    🖼️  Adding {len(visual_elements)} visual elements...")
        for visual_element in visual_elements:
            # Skip the full background we already added
            if visual_element['tag'] != 'full_background':
                self.add_visual_element_to_slide(slide, visual_element)
        
        # Step 3: Create editable text boxes on top of everything
        print(f"    ✍️  Adding {len(text_elements)} editable text elements...")
        for text_element in text_elements:
            self.create_text_box(slide, text_element)
        
        print(f"    🎉 Slide {slide_num} built successfully!")
    
    async def convert_slide_perfect(self, browser, slide_info: Dict, presentation, temp_dir: Path) -> None:
        """Convert a single HTML slide using perfect 1:1 approach."""
        html_path = slide_info['path']
        slide_num = slide_info['number']
        
        print(f"Converting slide {slide_num}: {slide_info['title']} (Perfect 1:1 Mode)")
        
        # Add blank slide
        blank_slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(blank_slide_layout)
        
        # Step 1: Extract individual visual elements (icons, graphics, borders, etc.)
        print("  🎨 Extracting individual visual elements...")
        visual_elements = await self.extract_visual_elements(browser, html_path, temp_dir)
        
        # Step 2: Capture clean background with visual elements temporarily hidden
        print("  🎨 Capturing clean background (visual elements hidden)...")
        background_path = await self.capture_clean_background(browser, html_path, temp_dir, visual_elements)
        
        # Step 3: Add the clean background as the base layer
        print("  🖼️  Adding clean background as base layer...")
        background_element = {
            'type': 'background',
            'x': 0,
            'y': 0,
            'width': 1920,
            'height': 1080,
            'tag': 'clean_background',
            'image_path': background_path,
            'depth': -1
        }
        self.add_visual_element_to_slide(slide, background_element)
        
        # Step 4: Add individual visual elements on top of background
        print("  🖼️  Adding individual visual elements...")
        for visual_element in visual_elements:
            # Skip the full background we already added
            if visual_element['tag'] != 'full_background':
                self.add_visual_element_to_slide(slide, visual_element)
        
        # Step 5: Extract and add editable text elements
        print("  📝 Extracting editable text elements...")
        text_elements = await self.extract_text_elements(browser, html_path)
        
        # Step 6: Create editable text boxes on top of everything
        print("  ✍️  Adding editable text overlays...")
        for text_element in text_elements:
            self.create_text_box(slide, text_element)
        
        print(f"  🎉 Slide {slide_num}: Perfect background + {len(visual_elements)} visual elements + {len(text_elements)} editable text elements")
    
    async def convert_to_pptx_perfect(self, store_locally: bool = True) -> tuple:
        """Main perfect conversion method with concurrent processing."""
        print("🎯 Starting PERFECT 1:1 HTML to PPTX conversion...")
        print("📋 Method: Individual visual elements + Editable text overlay")
        print("🔄 Processing: Concurrent slide analysis + Sequential PPTX building")
        print("=" * 80)
        
        # Load metadata
        self.load_metadata()
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Launch browser for processing
            async with async_playwright() as p:
                print("🌐 Launching browser for concurrent processing...")
                browser = await p.chromium.launch(
                    headless=True,
                    args=[
                        '--no-sandbox',
                        '--disable-setuid-sandbox',
                        '--disable-dev-shm-usage',
                        '--disable-gpu',
                        '--force-device-scale-factor=1',
                        '--disable-background-timer-throttling',
                        '--disable-backgrounding-occluded-windows',
                        '--disable-renderer-backgrounding',
                        '--disable-features=VizDisplayCompositor'
                    ]
                )
                
                try:
                    # STEP 1: Concurrent slide analysis (extract visual elements, capture backgrounds, extract text)
                    print(f"📄 Processing {len(self.slides_info)} slides concurrently...")
                    
                    async def analyze_slide(slide_info: Dict) -> Dict:
                        """Analyze a single slide concurrently - extract all data needed for PPTX creation."""
                        html_path = slide_info['path']
                        slide_num = slide_info['number']
                        
                        print(f"  🔍 Analyzing slide {slide_num}: {slide_info['title']}")
                        
                        # Extract visual elements
                        visual_elements = await self.extract_visual_elements(browser, html_path, temp_path)
                        
                        # Capture clean background
                        background_path = await self.capture_clean_background(browser, html_path, temp_path, visual_elements)
                        
                        # Extract text elements
                        text_elements = await self.extract_text_elements(browser, html_path)
                        
                        return {
                            'slide_info': slide_info,
                            'visual_elements': visual_elements,
                            'background_path': background_path,
                            'text_elements': text_elements
                        }
                    
                    # Process all slides concurrently using asyncio.gather
                    analysis_tasks = [
                        analyze_slide(slide_info) for slide_info in self.slides_info
                    ]
                    
                    # Wait for all slides to be analyzed concurrently
                    slide_analyses = await asyncio.gather(*analysis_tasks)
                    
                    print(f"✅ All {len(slide_analyses)} slides analyzed concurrently!")
                    
                finally:
                    await browser.close()
            
            # STEP 2: Sequential PPTX building (using the analyzed data)
            print("🏗️  Building PowerPoint presentation from analyzed data...")
            
            # Create new PowerPoint presentation
            presentation = Presentation()
            
            # Set slide dimensions to 1920x1080 (16:9)
            presentation.slide_width = Inches(20)  # 1920px at 96 DPI
            presentation.slide_height = Inches(11.25)  # 1080px at 96 DPI
            
            # Remove default slide
            if len(presentation.slides) > 0:
                xml_slides = presentation.slides._sldIdLst
                xml_slides.remove(xml_slides[0])
            
            # Build slides using the analyzed data (this is fast since all heavy work is done)
            for slide_analysis in slide_analyses:
                await self.build_slide_from_analysis(presentation, slide_analysis, temp_path)
            
            # Save PowerPoint presentation to temp (INSIDE the temp directory context)
            presentation_name = self.metadata.get('presentation_name', 'presentation')
            temp_output_path = temp_path / f"{presentation_name}.pptx"
            
            print(f"    💾 Saving presentation to temp: {temp_output_path}")
            print(f"    📁 Temp directory contents before save: {list(temp_path.glob('*'))}")
            
            presentation.save(str(temp_output_path))
            
            print(f"    ✅ Presentation saved to temp")
            print(f"    📁 Temp directory contents after save: {list(temp_path.glob('*'))}")
            print(f"    🔍 Checking if temp file exists: {temp_output_path.exists()}")
            if temp_output_path.exists():
                print(f"    📏 Temp file size: {temp_output_path.stat().st_size} bytes")
            else:
                print(f"    ❌ Temp file does not exist!")
            
            if store_locally:
                # Store in the static files directory for URL serving
                timestamp = int(asyncio.get_event_loop().time())
                filename = f"{presentation_name}_{timestamp}.pptx"
                final_output = output_dir / filename
                final_output.parent.mkdir(exist_ok=True)
                
                # Copy from temp to final location
                import shutil
                shutil.copy2(temp_output_path, final_output)
                
                print(f"\n🎉 PERFECT 1:1 PPTX created successfully: {final_output}")
                print(f"📊 Total slides: {len(presentation.slides)}")
                print(f"✨ Perfect visual fidelity + Fully editable text!")
                print(f"🚀 Concurrent processing completed!")
                
                return final_output, len(presentation.slides)
            else:
                # For direct download, read content from temp file before cleanup
                print(f"    📖 Reading file content for download...")
                try:
                    with open(temp_output_path, 'rb') as f:
                        pptx_content = f.read()
                    print(f"    ✅ File content read successfully: {len(pptx_content)} bytes")
                except Exception as e:
                    print(f"    ❌ Error reading file: {e}")
                    print(f"    📁 Final temp directory contents: {list(temp_path.glob('*'))}")
                    raise
                
                print(f"\n🎉 PERFECT 1:1 PPTX content prepared for download")
                print(f"📊 Total slides: {len(presentation.slides)}")
                print(f"✨ Perfect visual fidelity + Fully editable text!")
                print(f"🚀 Concurrent processing completed!")
                
                return pptx_content, len(presentation.slides), presentation_name


@router.post("/convert-to-pptx")
async def convert_presentation_to_pptx(request: ConvertRequest):
    """
    Convert HTML presentation to PPTX with perfect 1:1 fidelity and async processing.
    
    Takes a presentation folder path and returns either:
    - PPTX file directly (if download=true) - uses presentation name as filename
    - JSON response with download URL (if download=false, default)
    """
    try:
        print(f"📥 Received PPTX conversion request for: {request.presentation_path}")
        
        # Create converter
        converter = PerfectHTMLToPPTXConverter(request.presentation_path)
        
        # If download is requested, don't store locally and return file directly
        if request.download:
            pptx_content, total_slides, presentation_name = await converter.convert_to_pptx_perfect(store_locally=False)
            
            print(f"✨ Direct download conversion completed for: {presentation_name}")
            
            return Response(
                content=pptx_content,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename=\"{presentation_name}.pptx\""}
            )
        
        # Otherwise, store locally and return JSON with download URL
        pptx_path, total_slides = await converter.convert_to_pptx_perfect(store_locally=True)
        
        print(f"✨ Conversion completed: {pptx_path}")
        
        pptx_url = f"/downloads/{pptx_path.name}"
        
        return ConvertResponse(
            success=True,
            message=f"PPTX generated successfully with {total_slides} slides",
            pptx_url=pptx_url,
            filename=pptx_path.name,
            total_slides=total_slides
        )
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"❌ PPTX conversion error: {e}")
        raise HTTPException(status_code=500, detail=f"PPTX conversion failed: {str(e)}")


@router.get("/health")
async def pptx_health_check():
    """PPTX service health check endpoint."""
    return {"status": "healthy", "service": "HTML to PPTX Converter"}
