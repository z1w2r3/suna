#!/usr/bin/env python3
"""
Optimized FastAPI HTML Presentation to PPTX Converter Router

Provides PPTX conversion endpoints with improved memory management and batch processing.
Uses smart batching and resource cleanup to handle large presentations efficiently.
"""

import json
import asyncio
import os
from pathlib import Path
from typing import Dict, List, Optional
import tempfile
import shutil
from dataclasses import dataclass

from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel, Field

try:
    from playwright.async_api import async_playwright
except ImportError:
    raise ImportError("Playwright is not installed. Please install it with: pip install playwright")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    raise ImportError(f"python-pptx is not installed. Please install it with: pip install python-pptx. Error: {e}")


# Create router
router = APIRouter(prefix="/presentation", tags=["pptx-conversion"])

# Create output directory for generated PPTXs
output_dir = Path("generated_pptx")
output_dir.mkdir(exist_ok=True)


class ConvertRequest(BaseModel):
    presentation_path: str = Field(..., description="Path to the presentation folder containing metadata.json")
    download: bool = Field(False, description="If true, returns the PPTX file directly. If false, returns JSON with download URL.")


class ConvertResponse(BaseModel):
    success: bool
    message: str
    pptx_url: str
    filename: str
    total_slides: int


@dataclass
class TextElement:
    """Text element information for editable text boxes with enhanced styling"""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_family: str
    font_size: float
    font_weight: str
    color: str
    text_align: str
    line_height: float
    tag: str
    depth: int = 0
    style: Dict[str, str] = None


class CSSParser:
    """Parse CSS styles and convert values to appropriate units"""
    
    @staticmethod
    def parse_color(color_str: str) -> tuple[int, int, int]:
        """Parse CSS color string to RGB tuple"""
        if not color_str:
            return (0, 0, 0)
        
        color_str = color_str.strip().lower()
        
        # Handle hex colors
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            try:
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            except ValueError:
                return (0, 0, 0)
        
        # Handle rgb() colors
        import re
        rgb_match = re.match(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
        if rgb_match:
            return tuple(int(x) for x in rgb_match.groups())
        
        # Handle rgba() colors (ignore alpha for now)
        rgba_match = re.match(r'rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_str)
        if rgba_match:
            return tuple(int(x) for x in rgba_match.groups())
        
        # Named colors
        named_colors = {
            'black': (0, 0, 0), 'white': (255, 255, 255), 'red': (255, 0, 0),
            'green': (0, 128, 0), 'blue': (0, 0, 255), 'yellow': (255, 255, 0),
            'cyan': (0, 255, 255), 'magenta': (255, 0, 255), 'gray': (128, 128, 128),
            'grey': (128, 128, 128), 'orange': (255, 165, 0), 'purple': (128, 0, 128)
        }
        
        return named_colors.get(color_str, (0, 0, 0))
    
    @staticmethod
    def parse_font_weight(weight_str: str) -> bool:
        """Parse font weight to bold boolean"""
        if not weight_str:
            return False
        
        weight_str = weight_str.strip().lower()
        bold_weights = ['bold', 'bolder', '700', '800', '900']
        
        return weight_str in bold_weights or (weight_str.isdigit() and int(weight_str) >= 700)


class OptimizedHTMLToPPTXConverter:
    def __init__(self, presentation_dir: str):
        """Initialize the optimized converter."""
        self.presentation_dir = Path(presentation_dir).resolve()
        self.metadata_path = self.presentation_dir / "metadata.json"
        self.metadata = None
        self.slides_info = []
        
        # Validate inputs
        if not self.presentation_dir.exists():
            raise FileNotFoundError(f"Presentation directory not found: {self.presentation_dir}")
        
        if not self.metadata_path.exists():
            raise FileNotFoundError(f"metadata.json not found in: {self.presentation_dir}")
    
    def load_metadata(self) -> Dict:
        """Load and parse metadata.json"""
        try:
            with open(self.metadata_path, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
            
            # Extract slide information and sort by slide number
            slides = self.metadata.get('slides', {})
            self.slides_info = []
            
            for slide_num, slide_data in slides.items():
                filename = slide_data.get('filename')
                file_path = slide_data.get('file_path')
                title = slide_data.get('title', f'Slide {slide_num}')
                
                if file_path:
                    # Handle both absolute and relative paths
                    if Path(file_path).is_absolute():
                        html_path = Path(file_path)
                    else:
                        html_path = Path(f"/workspace/{file_path}")
                    
                    # Verify the path exists
                    if html_path.exists():
                        self.slides_info.append({
                            'number': int(slide_num),
                            'title': title,
                            'filename': filename,
                            'path': html_path
                        })
            
            # Sort slides by number
            self.slides_info.sort(key=lambda x: x['number'])
            
            if not self.slides_info:
                raise ValueError("No valid slides found in metadata.json")
            
            return self.metadata
            
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in metadata.json: {e}")
        except Exception as e:
            raise ValueError(f"Error loading metadata: {e}")
    
    async def extract_visual_elements(self, page, html_path: Path, temp_dir: Path) -> List[Dict]:
        """Extract all visual elements (non-text) as individual images with positioning."""
        visual_elements = []
        
        try:
            # Set viewport and load HTML
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Use file:// URL instead of set_content to preserve relative paths
            file_url = f"file://{html_path.resolve()}"
            await page.goto(file_url, wait_until="networkidle", timeout=25000)
            await page.wait_for_timeout(1000)
            
            def handle_console(msg):
                print(f"BROWSER CONSOLE: {msg.text}")

            page.on("console", handle_console)
            
            # Step 1: First extract icons BEFORE making text transparent
            icon_data = await page.evaluate(r"""
                () => {
                    function isIcon(element) {
                        // Check for Font Awesome icons
                        if (element.classList.contains('fas') || 
                            element.classList.contains('far') || 
                            element.classList.contains('fab') || 
                            element.classList.contains('fa')) {
                            return true;
                        }
                        
                        return false;
                    }
                    
                    function extractIconElements(element, depth = 0) {
                        if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
                        
                        const computed = window.getComputedStyle(element);
                        const rect = element.getBoundingClientRect();
                        
                        // Skip if no dimensions or hidden
                        if (rect.width === 0 || rect.height === 0) return [];
                        if (computed.display === 'none' || computed.visibility === 'hidden') return [];
                        
                        const results = [];
                        
                        // Check if this element is an icon
                        if (isIcon(element)) {
                            console.log('ICON ELEMENT FOUND:', element.tagName, element.textContent?.trim().substring(0, 30));
                            
                            results.push({
                                type: 'icon',
                                x: Math.round(rect.left * 100) / 100,
                                y: Math.round(rect.top * 100) / 100,
                                width: Math.round(rect.width * 100) / 100,
                                height: Math.round(rect.height * 100) / 100,
                                tag: element.tagName.toLowerCase(),
                                className: element.className,
                                id: element.id,
                                depth: depth,
                                // Add debug info
                                debugInfo: {
                                    textContent: element.textContent?.trim().substring(0, 50) || 'No text'
                                }
                            });
                        }
                        
                        // Process children for nested icons
                        Array.from(element.children).forEach(child => {
                            results.push(...extractIconElements(child, depth + 1));
                        });
                        
                        return results;
                    }
                    
                    // Execute icon extraction
                    const allIconElements = extractIconElements(document.body);
                    
                    return allIconElements;
                }
            """)
            
            # Capture icons as visual elements
            print(f"🔍 Total icon elements found: {len(icon_data) if icon_data else 0}")
            icon_visual_elements = []
            
            for i, data in enumerate(icon_data or []):
                if data and data['type'] == 'icon':
                    try:
                        # Ensure coordinates are within viewport bounds
                        x = max(0, min(data['x'], 1920))
                        y = max(0, min(data['y'], 1080))
                        width = min(data['width'], 1920 - x)
                        height = min(data['height'], 1080 - y)
                        
                        # Skip if area is too small
                        if width < 5 or height < 5:
                            continue
                        
                        # Capture the icon
                        element_path = temp_dir / f"icon_element_{html_path.stem}_{i:03d}.png"
                        await page.screenshot(
                            path=str(element_path),
                            full_page=False,
                            clip={"x": x, "y": y, "width": width, "height": height}
                        )
                        
                        icon_element = {
                            'type': 'visual',  # Treat as visual element for consistency
                            'x': data['x'],
                            'y': data['y'],
                            'width': data['width'],
                            'height': data['height'],
                            'tag': data['tag'],
                            'image_path': element_path,
                            'depth': data['depth'],
                            'isIcon': True
                        }
                        
                        icon_visual_elements.append(icon_element)
                        
                    except Exception as e:
                        print(f"Failed to capture icon element {i}: {e}")
            
            # Step 2: Now make all text transparent while preserving visual styling
            await page.evaluate(r"""
                () => {
                    function makeTextTransparent(element) {
                        if (element.nodeType === Node.ELEMENT_NODE) {
                            const computed = window.getComputedStyle(element);
                            
                            // If element has text content, make it transparent
                            if (element.textContent && element.textContent.trim()) {
                                element.style.color = 'transparent';
                                element.style.textShadow = 'none';
                                element.style.webkitTextStroke = 'none';
                                element.style.webkitTextFillColor = 'transparent';
                            }
                            
                            // Process children
                            Array.from(element.children).forEach(makeTextTransparent);
                        }
                    }
                    
                    makeTextTransparent(document.body);
                }
            """)
            
            await page.wait_for_timeout(500)
            
            # Step 3: Extract other visual elements by depth and visual properties
            visual_data = await page.evaluate(r"""
    () => {
        function hasActualVisualContent(element, computed) {
            
            // Always include explicit visual elements
            if (['IMG', 'SVG', 'CANVAS', 'VIDEO', 'IFRAME'].includes(element.tagName)) {
                return true;
            }
            
            // 2. Check for actual background images and gradients (not just 'none')
            if (computed.backgroundImage && computed.backgroundImage !== 'none') {
                return true;
            }
            
            // 3. Check for meaningful background colors (not just transparent)
            const bgColor = computed.backgroundColor;
            if (bgColor && 
                bgColor !== 'rgba(0, 0, 0, 0)' && 
                bgColor !== 'transparent' && 
                bgColor !== 'inherit' &&
                bgColor !== 'initial' &&
                bgColor !== 'unset') {
                return true;
            }
            
            // 4. Check for borders (but ignore default/none)
            if (computed.borderStyle && 
                computed.borderStyle !== 'none' && 
                computed.borderStyle !== 'initial' &&
                computed.borderWidth && 
                computed.borderWidth !== '0px') {
                return true;
            }
            
            // 5. Check for box shadows
            if (computed.boxShadow && 
                computed.boxShadow !== 'none' && 
                computed.boxShadow !== 'initial') {
                return true;
            }
            
            // 6. Check for gradients in background (comprehensive check)
            if (computed.background && 
                (computed.background.includes('gradient') || 
                 computed.background.includes('url('))) {
                return true;
            }
            
            // Also check backgroundImage for gradients (browsers often store gradients here)
            if (computed.backgroundImage && 
                (computed.backgroundImage.includes('gradient') ||
                 computed.backgroundImage.includes('url('))) {
                return true;
            }
            
            return false;
        }
                          
        function shouldSkipElement(element, computed) {
            // Skip text-only elements
            const textOnlyTags = ['H1', 'H2', 'H3', 'H4', 'H5', 'H6', 'P', 'A', 'SPAN', 
                                'STRONG', 'EM', 'U', 'BUTTON', 'LABEL', 'SMALL', 'CODE'];
            if (textOnlyTags.includes(element.tagName)) {
                return true;
            }
            
            // Skip hidden elements
            if (computed.display === 'none' || computed.visibility === 'hidden') {
                return true;
            }
            
            return false;
        }
                                    
        function extractVisualElements(element, depth = 0) {
            if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
            
            const computed = window.getComputedStyle(element);
            const rect = element.getBoundingClientRect();
            
            // Skip if no dimensions
            if (rect.width === 0 || rect.height === 0) {
                return [];
            }
            
            // Skip if should be filtered out
            if (shouldSkipElement(element, computed)) {
                return [];
            }
            
            const results = [];
            
            // Check if element has actual visual content
            if (hasActualVisualContent(element, computed)) {
                console.log('VISUAL ELEMENT FOUND:', element.tagName, element.textContent?.trim().substring(0, 30));
                console.log('backgroundImage:', computed.backgroundImage);
                console.log('backgroundColor:', computed.backgroundColor);
                console.log('borderStyle:', computed.borderStyle);
                
                // Skip very large elements that are likely backgrounds
                const isLikelyBackground = rect.width > 1700;
                
                if (!isLikelyBackground) {
                    // Assign unique capture ID to element
                    const captureId = 'visual-capture-' + Date.now() + '-' + Math.random().toString(36).substr(2, 9);
                    element.setAttribute('data-capture-id', captureId);
                    
                    results.push({
                        type: 'visual',
                        captureId: captureId,
                        x: Math.round(rect.left * 100) / 100,
                        y: Math.round(rect.top * 100) / 100,
                        width: Math.round(rect.width * 100) / 100,
                        height: Math.round(rect.height * 100) / 100,
                        tag: element.tagName.toLowerCase(),
                        className: element.className,
                        id: element.id,
                        depth: depth,
                        hasBackground: computed.backgroundImage !== 'none' || 
                                      (computed.backgroundColor !== 'rgba(0, 0, 0, 0)' && 
                                       computed.backgroundColor !== 'transparent'),
                        hasBorder: computed.borderStyle !== 'none',
                        hasShadow: computed.boxShadow !== 'none',
                        // Add debug info
                        debugInfo: {
                            backgroundColor: computed.backgroundColor,
                            backgroundImage: computed.backgroundImage,
                            borderStyle: computed.borderStyle,
                            textContent: element.textContent?.trim().substring(0, 50) || 'No text'
                        }
                    });
                }
            }
            
            // Process children for nested elements
            Array.from(element.children).forEach(child => {
                results.push(...extractVisualElements(child, depth + 1));
            });
            
            return results;
        }
        
        // Actually execute the extraction and return results
        const allVisualElements = extractVisualElements(document.body);
        
        // Sort by depth (background elements first) then by position
        allVisualElements.sort((a, b) => {
            if (a.depth !== b.depth) return a.depth - b.depth;
            if (Math.abs(a.y - b.y) < 5) return a.x - b.x;
            return a.y - b.y;
        });
        
        return allVisualElements;
    }
""")
            
            # Debug: Log what we found in Python console
            print(f"🔍 Total visual elements found: {len(visual_data) if visual_data else 0}")
            if visual_data:
                for i, data in enumerate(visual_data):  # Show first 5
                    print(f"🔍 Element {i}: {data['tag']} - {data.get('className', 'No class')}")
                    print(f"   Position: ({data['x']}, {data['y']}) {data['width']}x{data['height']}")
                    print(f"   Has background: {data.get('hasBackground', False)}")
                    print(f"   Has border: {data.get('hasBorder', False)}")
                    print(f"   Has shadow: {data.get('hasShadow', False)}")
                    print("--- 🔍", data)
            
            # Step 3: Capture each visual element as an image
            for i, data in enumerate(visual_data or []):
                if data and data['type'] == 'visual':
                    try:
                        # Ensure coordinates are within viewport bounds
                        x = max(0, min(data['x'], 1920))
                        y = max(0, min(data['y'], 1080))
                        width = min(data['width'], 1920 - x)
                        height = min(data['height'], 1920 - y)
                        
                        # Skip if area is too small
                        if width < 5 or height < 5:
                            continue
                        
                        # Clone element and capture in isolation
                        clone_result = await page.evaluate(r"""
                            (elementData) => {
                                // Find the element by its unique capture ID
                                const targetElement = document.querySelector(`[data-capture-id="${elementData.captureId}"]`);
                                
                                if (!targetElement) {
                                    return { success: false, error: `Visual element not found with ID: ${elementData.captureId}` };
                                }
                                
                                try {
                                    // Clone the element deeply
                                    const clonedElement = targetElement.cloneNode(true);
                                    
                                    // Copy all computed styles from original to cloned element and its descendants
                                    function copyComputedStyles(original, cloned, isRoot = true) {
                                        const computedStyle = window.getComputedStyle(original);
                                        
                                        // Create a clean style string with all computed properties
                                        let styleStr = '';
                                        for (let prop of computedStyle) {
                                            let value = computedStyle.getPropertyValue(prop);
                                            styleStr += `${prop}: ${value}; `;
                                        }
                                        cloned.style.cssText = styleStr;
                                        
                                        // For child elements (not root), make the entire element transparent
                                        if (!isRoot) {
                                            cloned.style.opacity = '0';
                                        }
                                        
                                        // Recursively copy styles for children (mark as non-root)
                                        for (let i = 0; i < original.children.length && i < cloned.children.length; i++) {
                                            copyComputedStyles(original.children[i], cloned.children[i], false);
                                        }
                                    }
                                    
                                    copyComputedStyles(targetElement, clonedElement);
                                    
                                    // Create clean container for isolation
                                    const cleanContainer = document.createElement('div');
                                    cleanContainer.id = 'clone-capture-container-' + Date.now();
                                    cleanContainer.style.cssText = `
                                        position: fixed;
                                        top: 0;
                                        left: 0;
                                        width: ${elementData.width}px;
                                        height: ${elementData.height}px;
                                        background: transparent;
                                        z-index: 999999;
                                        padding: 0;
                                        margin: 0;
                                        border: none;
                                        overflow: hidden;
                                    `;
                                    
                                    // Position cloned element at (0,0) within container
                                    clonedElement.style.position = 'absolute';
                                    clonedElement.style.top = '0px';
                                    clonedElement.style.left = '0px';
                                    clonedElement.style.margin = '0';
                                    clonedElement.style.transform = 'none';
                                    
                                    // Add to DOM for capture
                                    cleanContainer.appendChild(clonedElement);
                                    document.body.appendChild(cleanContainer);
                                    
                                    return { 
                                        success: true, 
                                        containerId: cleanContainer.id,
                                        containerRect: {
                                            x: 0, 
                                            y: 0, 
                                            width: elementData.width, 
                                            height: elementData.height
                                        }
                                    };
                                    
                                } catch (error) {
                                    return { success: false, error: error.message };
                                }
                            }
                        """, data)
                        
                        if not clone_result.get('success'):
                            print(f"Failed to clone element: {clone_result.get('error')}")
                            continue
                        
                        # Wait a moment for styles to apply
                        await page.wait_for_timeout(100)
                        
                        # Hide main content and ensure transparent background for perfect isolation
                        await page.evaluate(f"""
                            () => {{
                                // Store original backgrounds
                                window.originalHtmlBg = document.documentElement.style.background || '';
                                window.originalBodyBg = document.body.style.background || '';
                                
                                // Set transparent backgrounds
                                document.documentElement.style.background = 'transparent';
                                document.body.style.background = 'transparent';
                                
                                // Hide all main content
                                document.body.style.visibility = 'hidden';
                                
                                // Make sure our container stays visible
                                const container = document.getElementById('{clone_result['containerId']}');
                                if (container) {{
                                    container.style.visibility = 'visible';
                                }}
                            }}
                        """)
                        
                        # Capture the isolated cloned element with transparency
                        element_path = temp_dir / f"visual_element_{html_path.stem}_{i:03d}.png"
                        container_rect = clone_result['containerRect']
                        
                        await page.screenshot(
                            path=str(element_path),
                            full_page=False,
                            omit_background=True,  # 🚀 Preserves transparency!
                            clip={
                                "x": container_rect['x'], 
                                "y": container_rect['y'], 
                                "width": container_rect['width'], 
                                "height": container_rect['height']
                            }
                        )
                        
                        # Restore backgrounds, visibility and clean up
                        await page.evaluate(f"""
                            () => {{
                                // Restore original backgrounds
                                document.documentElement.style.background = window.originalHtmlBg || '';
                                document.body.style.background = window.originalBodyBg || '';
                                
                                // Restore main content visibility
                                document.body.style.visibility = 'visible';
                                
                                // Remove the cloned container
                                const container = document.getElementById('{clone_result['containerId']}');
                                if (container) {{
                                    container.remove();
                                }}
                            }}
                        """)
                        
                        visual_element = {
                            'type': 'visual',
                            'x': data['x'],
                            'y': data['y'],
                            'width': data['width'],
                            'height': data['height'],
                            'tag': data['tag'],
                            'image_path': element_path,
                            'depth': data['depth'],
                            'hasBackground': data.get('hasBackground', False),
                            'hasBorder': data.get('hasBorder', False),
                            'hasShadow': data.get('hasShadow', False)
                        }
                        
                        visual_elements.append(visual_element)
                        
                    except Exception as e:
                        print(f"Failed to capture visual element {i}: {e}")
                        # Ensure we restore backgrounds and visibility even if capture fails
                        try:
                            await page.evaluate("""
                                () => {
                                    // Restore original backgrounds
                                    document.documentElement.style.background = window.originalHtmlBg || '';
                                    document.body.style.background = window.originalBodyBg || '';
                                    
                                    // Always restore main content visibility
                                    document.body.style.visibility = 'visible';
                                    
                                    // Clean up any leftover containers
                                    const containers = document.querySelectorAll('[id^="clone-capture-container-"]');
                                    containers.forEach(container => container.remove());
                                }
                            """)
                        except:
                            pass
                        continue
            
            # Add the previously captured icon elements to our visual elements collection
            if icon_visual_elements:
                print(f"Adding {len(icon_visual_elements)} icon elements to visual elements")
                visual_elements.extend(icon_visual_elements)
            
            # Final cleanup - remove capture IDs and ensure everything is restored
            try:
                await page.evaluate(r"""
                    () => {
                        // Restore original backgrounds (final safety check)
                        if (window.originalHtmlBg !== undefined) {
                            document.documentElement.style.background = window.originalHtmlBg || '';
                        }
                        if (window.originalBodyBg !== undefined) {
                            document.body.style.background = window.originalBodyBg || '';
                        }
                        
                        // Ensure visibility is restored
                        document.body.style.visibility = 'visible';
                        
                        // Clean up capture IDs
                        const elementsWithCaptureId = document.querySelectorAll('[data-capture-id]');
                        elementsWithCaptureId.forEach(el => {
                            el.removeAttribute('data-capture-id');
                        });
                        
                        // Clean up any leftover containers
                        const containers = document.querySelectorAll('[id^="clone-capture-container-"]');
                        containers.forEach(container => container.remove());
                    }
                """)
                print(f"🧹 Final cleanup completed: backgrounds restored, capture IDs removed")
            except Exception as e:
                print(f"Warning: Failed to perform final cleanup: {e}")
            
            return visual_elements
            
        except Exception as e:
            print(f"Visual element extraction failed: {e}")
            # Emergency cleanup in case of failure
            try:
                await page.evaluate(r"""
                    () => {
                        // Restore original backgrounds
                        if (window.originalHtmlBg !== undefined) {
                            document.documentElement.style.background = window.originalHtmlBg || '';
                        }
                        if (window.originalBodyBg !== undefined) {
                            document.body.style.background = window.originalBodyBg || '';
                        }
                        
                        // Ensure visibility is restored
                        document.body.style.visibility = 'visible';
                        
                        // Clean up everything
                        const containers = document.querySelectorAll('[id^="clone-capture-container-"]');
                        containers.forEach(container => container.remove());
                        
                        const elementsWithCaptureId = document.querySelectorAll('[data-capture-id]');
                        elementsWithCaptureId.forEach(el => el.removeAttribute('data-capture-id'));
                    }
                """)
            except:
                pass
            return []

    async def capture_clean_background(self, page, html_path: Path, temp_dir: Path, visual_elements: List[Dict]) -> Path:
        """Capture the clean background with visual elements temporarily hidden."""
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1 for exact measurements
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Use file:// URL instead of set_content to preserve relative paths
            file_url = f"file://{html_path.resolve()}"
            await page.goto(file_url, wait_until="networkidle", timeout=25000)
            
            # Reduced wait time
            await page.wait_for_timeout(2000)
            
            # Make text completely invisible AND hide visual elements to get clean background
            await page.evaluate(r"""
                (visualElementsData) => {
                    // Function to make text completely invisible
                    function makeTextInvisible(element) {
                        if (element.nodeType === Node.TEXT_NODE) return;
                        if (element.nodeType === Node.ELEMENT_NODE) {
                            const computedStyle = window.getComputedStyle(element);
                            const hasText = element.textContent && element.textContent.trim();
                            if (hasText) {
                                // Make text completely invisible
                                element.style.color = 'transparent';
                                element.style.textShadow = 'none';
                                element.style.webkitTextStroke = 'none';
                                element.style.webkitTextFillColor = 'transparent';
                                
                                // Also hide any background text effects
                                if (computedStyle.webkitBackgroundClip === 'text') {
                                    element.style.background = 'transparent';
                                    element.style.webkitBackgroundClip = 'initial';
                                }
                            }
                            Array.from(element.children).forEach(makeTextInvisible);
                        }
                    }
                    
                    // Function to hide visual elements (but preserve main background)
                    function hideVisualElements() {
                        visualElementsData.forEach(visualData => {
                            // Skip very large elements that are likely the main background
                            if (visualData.width > 1700) {
                                return;
                            }
                            
                            // Find elements that match the visual element criteria
                            const elements = document.querySelectorAll('*');
                            for (let el of elements) {
                                const rect = el.getBoundingClientRect();
                                if (Math.abs(rect.left - visualData.x) < 5 && 
                                    Math.abs(rect.top - visualData.y) < 5 &&
                                    Math.abs(rect.width - visualData.width) < 5 &&
                                    Math.abs(rect.height - visualData.height) < 5) {
                                    
                                    // Hide this visual element
                                    el.style.visibility = 'hidden';
                                    break;
                                }
                            }
                        });
                    }
                    
                    // Apply both transformations
                    makeTextInvisible(document.body);
                    hideVisualElements();
                }
            """, visual_elements or [])
            
            # Reduced wait time
            await page.wait_for_timeout(500)
            
            # Take clean background screenshot
            background_path = temp_dir / f"clean_background_{html_path.stem}.png"
            await page.screenshot(
                path=str(background_path),
                full_page=False,
                clip={"x": 0, "y": 0, "width": 1920, "height": 1080}
            )
            
            return background_path
            
        except Exception as e:
            # Create a simple white background as fallback
            from PIL import Image
            background_path = temp_dir / f"clean_background_{html_path.stem}.png"
            blank_bg = Image.new('RGB', (1920, 1080), color='white')
            blank_bg.save(background_path)
            return background_path
    
    async def extract_text_elements(self, page, html_path: Path) -> List[TextElement]:
        """Extract all text elements with precise positioning for editable text boxes."""
        text_elements = []
        
        try:
            # Set exact viewport dimensions
            await page.set_viewport_size({"width": 1920, "height": 1080})
            await page.emulate_media(media='screen')
            
            # Force device pixel ratio to 1 for exact measurements
            await page.evaluate(r"""
                () => {
                    Object.defineProperty(window, 'devicePixelRatio', {
                        get: () => 1
                    });
                }
            """)
            
            # Use file:// URL instead of set_content to preserve relative paths
            file_url = f"file://{html_path.resolve()}"
            await page.goto(file_url, wait_until="networkidle", timeout=25000)
            
            # Reduced wait time
            await page.wait_for_timeout(2000)
            
            # Extract all text elements with precise positioning and styling
            # // Enhanced text extraction JavaScript to include in your page.evaluate()
            text_data = await page.evaluate(r"""
                () => {
                    function extractTextFromElement(element) {
                        if (!element || element.nodeType !== Node.ELEMENT_NODE) return [];
                        
                        const computedStyle = window.getComputedStyle(element);
                        const rect = element.getBoundingClientRect();
                        
                        // Skip hidden elements
                        if (computedStyle.display === 'none' || computedStyle.visibility === 'hidden') return [];
                        if (rect.width === 0 || rect.height === 0) return [];
                        
                        const results = [];
                        
                        // Get direct text content (not from children)
                        let directText = '';
                        for (let node of element.childNodes) {
                            if (node.nodeType === Node.TEXT_NODE) {
                                directText += node.textContent;
                            }
                        }
                        directText = directText.trim();
                        
                        // If this element has direct text content, extract it with styling
                        if (directText && directText.length > 0) {
                            const fontSizeMatch = computedStyle.fontSize.match(/([0-9.]+)px/);
                            const actualFontSize = fontSizeMatch ? parseFloat(fontSizeMatch[1]) : 16;
                            
                            // Get the actual text node position, not the container
                            const textNodes = [];
                            for (let node of element.childNodes) {
                                if (node.nodeType === Node.TEXT_NODE && node.textContent.trim()) {
                                    const range = document.createRange();
                                    range.selectNodeContents(node);
                                    const textRect = range.getBoundingClientRect();
                                    textNodes.push({
                                        text: node.textContent.trim(),
                                        rect: textRect
                                    });
                                }
                            }
                            
                            // Enhanced styling extraction
                            const textStyle = {
                                fontFamily: computedStyle.fontFamily,
                                fontSize: computedStyle.fontSize,
                                fontWeight: computedStyle.fontWeight,
                                fontStyle: computedStyle.fontStyle,
                                color: computedStyle.color,
                                textAlign: computedStyle.textAlign,
                                lineHeight: computedStyle.lineHeight,
                                letterSpacing: computedStyle.letterSpacing,
                                textShadow: computedStyle.textShadow,
                                webkitTextStroke: computedStyle.webkitTextStroke,
                                webkitTextFillColor: computedStyle.webkitTextFillColor,
                                background: computedStyle.background,
                                backgroundImage: computedStyle.backgroundImage,
                                webkitBackgroundClip: computedStyle.webkitBackgroundClip,
                                textDecoration: computedStyle.textDecoration,
                                textTransform: computedStyle.textTransform
                            };
                            
                            // Use text node position if available, otherwise fall back to container
                            if (textNodes.length > 0) {
                                textNodes.forEach(textNode => {
                                    results.push({
                                        text: textNode.text,
                                        x: Math.round(textNode.rect.left * 100) / 100,
                                        y: Math.round(textNode.rect.top * 100) / 100,
                                        width: Math.round(textNode.rect.width * 100) / 100,
                                        height: Math.round(textNode.rect.height * 100) / 100,
                                        actualFontSizePx: actualFontSize,
                                        tag: element.tagName.toLowerCase(),
                                        style: textStyle
                                    });
                                });
                            } else {
                                // Fallback to container position
                                results.push({
                                    text: directText,
                                    x: Math.round(rect.left * 100) / 100,
                                    y: Math.round(rect.top * 100) / 100,
                                    width: Math.round(rect.width * 100) / 100,
                                    height: Math.round(rect.height * 100) / 100,
                                    actualFontSizePx: actualFontSize,
                                    tag: element.tagName.toLowerCase(),
                                    style: textStyle
                                });
                            }
                        }
                        
                        // Process children for nested text
                        Array.from(element.children).forEach(child => {
                            results.push(...extractTextFromElement(child));
                        });
                        
                        return results;
                    }
                    
                    // Start extraction from body
                    const allTextElements = extractTextFromElement(document.body);
                    
                    // Sort by position (top to bottom, left to right)
                    allTextElements.sort((a, b) => {
                        if (Math.abs(a.y - b.y) < 5) {
                            return a.x - b.x;
                        }
                        return a.y - b.y;
                    });
                    
                    return allTextElements;
                }
            """)
            
            # Convert to TextElement objects with enhanced styling
            for data in text_data or []:
                if data and data['text']:
                    style = data.get('style', {})
                    
                    # Parse font family
                    font_family = style.get('fontFamily', 'Arial')
                    if font_family:
                        font_family = font_family.split(',')[0].strip().strip('"\'')
                        font_family_map = {
                            'roboto': 'Roboto', 'arial': 'Arial', 'helvetica': 'Helvetica',
                            'sans-serif': 'Arial', 'serif': 'Times New Roman', 'monospace': 'Courier New',
                            'jetbrains mono': 'Courier New', 'courier new': 'Courier New'
                        }
                        font_family = font_family_map.get(font_family.lower(), font_family)
                    else:
                        font_family = 'Arial'
                    
                    # Parse line height (same for both text and icons)
                    line_height = 1.2
                    line_height_str = style.get('lineHeight', 'normal')
                    if line_height_str and line_height_str != 'normal':
                        if line_height_str.endswith('px'):
                            px_value = float(line_height_str[:-2])
                            line_height = px_value / data['actualFontSizePx']
                        else:
                            try:
                                line_height = float(line_height_str)
                            except:
                                line_height = 1.2
                    
                    # Parse color - handle complex color scenarios
                    color = style.get('color', '#000000')
                    
                    # Handle gradient text (webkit background clip)
                    if style.get('webkitBackgroundClip') == 'text' and style.get('backgroundImage'):
                        bg_image = style.get('backgroundImage', '')
                        if 'linear-gradient' in bg_image:
                            import re
                            color_match = re.search(r'#[0-9a-fA-F]{6}', bg_image)
                            if color_match:
                                color = color_match.group(0)
                            else:
                                color = '#3B82F6'  # Default blue for gradients
                    
                    text_element = TextElement(
                        text=data['text'],
                        x=data['x'],
                        y=data['y'],
                        width=data['width'],
                        height=data['height'],
                        font_family=font_family,
                        font_size=data['actualFontSizePx'] * 0.75,  # Convert px to points
                        font_weight=style.get('fontWeight', 'normal'),
                        color=color,
                        text_align=style.get('textAlign', 'left'),
                        line_height=line_height,
                        tag=data['tag'],
                        style=style
                    )
                    
                    text_elements.append(text_element)
            
            return text_elements
            
        except Exception:
            return []
    
    def create_text_box(self, slide, text_element: TextElement) -> None:
        """Create an editable text box in PowerPoint with exact positioning and enhanced styling."""
        # Convert pixel coordinates to inches
        left = Inches(text_element.x / 96.0)
        top = Inches(text_element.y / 96.0)
        
        # Calculate proper width - ensure it's wide enough for the text content
        text_width = text_element.width
        
        # For large fonts, increase the width calculation significantly
        font_size_factor = max(1.0, text_element.font_size / 20.0)  # Scale factor based on font size
        
        if text_width < 100:  # If width is too small, estimate based on text length and font size
            # More generous estimate: each character needs more space for larger fonts
            chars_per_pixel = 8 / font_size_factor  # Fewer characters per pixel for larger fonts
            estimated_width = len(text_element.text) * chars_per_pixel
            text_width = max(text_width, estimated_width)
        
        # Add generous padding to prevent text from being cut off
        padding = max(20, text_element.font_size * 0.5)  # Padding scales with font size
        final_width = text_width + (padding * 2)
        
        width = Inches(final_width / 96.0)
        height = Inches(max(text_element.height, 10) / 96.0)
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Set text frame properties for exact positioning
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.word_wrap = True
        text_frame.auto_size = None
        
        # Add paragraph
        p = text_frame.paragraphs[0]
        p.text = text_element.text
        
        # Set text alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'centre': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY, 'start': PP_ALIGN.LEFT,
            'end': PP_ALIGN.RIGHT
        }
        p.alignment = alignment_map.get(text_element.text_align.lower(), PP_ALIGN.LEFT)
        
        # Set spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if hasattr(p, 'line_spacing'):
            p.line_spacing = text_element.line_height
        
        # Set font properties
        font = p.font
        font.name = text_element.font_family
        font.size = Pt(max(text_element.font_size, 8))
        font.bold = CSSParser.parse_font_weight(text_element.font_weight)
        
        # Enhanced styling from the style object
        if text_element.style:
            style = text_element.style
            
            # Handle letter spacing
            letter_spacing = style.get('letterSpacing', 'normal')
            if letter_spacing != 'normal' and letter_spacing.endswith('px'):
                try:
                    spacing_px = float(letter_spacing[:-2])
                    if hasattr(font, 'character_spacing'):
                        font.character_spacing = spacing_px
                except:
                    pass
            
            # Handle text transform
            text_transform = style.get('textTransform', 'none')
            if text_transform == 'uppercase':
                p.text = p.text.upper()
            elif text_transform == 'lowercase':
                p.text = p.text.lower()
            elif text_transform == 'capitalize':
                p.text = p.text.title()
        
        # Set font color with enhanced handling
        try:
            # Handle gradient text (webkit background clip)
            if text_element.style and text_element.style.get('webkitBackgroundClip') == 'text':
                r, g, b = CSSParser.parse_color(text_element.color)
                font.color.rgb = RGBColor(r, g, b)
            else:
                # Regular color handling
                r, g, b = CSSParser.parse_color(text_element.color)
                font.color.rgb = RGBColor(r, g, b)
        except Exception:
            font.color.rgb = RGBColor(0, 0, 0)
        
        # Make textbox transparent (no background, no border)
        textbox.fill.background()
        textbox.line.fill.background()
    
    def add_visual_element_to_slide(self, slide, visual_element: Dict) -> None:
        """Add a visual element as an image to the PowerPoint slide with exact positioning."""
        if visual_element['image_path'].exists():
            # Convert pixel coordinates to inches
            left = Inches(visual_element['x'] / 96.0)
            top = Inches(visual_element['y'] / 96.0)
            width = Inches(visual_element['width'] / 96.0)
            height = Inches(visual_element['height'] / 96.0)
            
            # Add the image to the slide
            picture = slide.shapes.add_picture(str(visual_element['image_path']), left, top, width, height)
            
            # Special handling for background elements
            if visual_element['tag'] == 'clean_background':
                picture.z_order = 0
    
    async def build_slide_from_analysis(self, presentation, slide_analysis: Dict, temp_dir: Path) -> None:
        """Build a PowerPoint slide from pre-analyzed data."""
        slide_info = slide_analysis['slide_info']
        visual_elements = slide_analysis['visual_elements']
        background_path = slide_analysis['background_path']
        text_elements = slide_analysis['text_elements']
        
        # Add blank slide
        blank_slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(blank_slide_layout)
        
        # Step 1: Add the clean background as the base layer
        if background_path and background_path.exists():
            background_element = {
                'type': 'background',
                'x': 0,
                'y': 0,
                'width': 1920,
                'height': 1080,
                'tag': 'clean_background',
                'image_path': background_path,
                'depth': -1
            }
            self.add_visual_element_to_slide(slide, background_element)
        
        # Step 2: Add individual visual elements on top of background
        if visual_elements:
            for visual_element in visual_elements:
                if visual_element['tag'] != 'full_background':
                    self.add_visual_element_to_slide(slide, visual_element)
        
        # Step 3: Create editable text boxes on top of everything
        if text_elements:
            for text_element in text_elements:
                try:
                    self.create_text_box(slide, text_element)
                except Exception:
                    pass
    
    async def convert_to_pptx(self, store_locally: bool = True) -> tuple:
        """Main conversion method - optimized and reliable."""
        # Load metadata
        self.load_metadata()
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Launch browser for processing
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True,
                    args=[
                        '--no-sandbox',
                        '--disable-setuid-sandbox',
                        '--disable-dev-shm-usage',
                        '--disable-gpu',
                        '--force-device-scale-factor=1',
                        '--disable-background-timer-throttling',
                        '--disable-backgrounding-occluded-windows',
                        '--disable-renderer-backgrounding',
                        '--disable-features=VizDisplayCompositor',
                        '--disable-extensions',
                        '--disable-plugins',
                        '--disable-web-security',
                        '--disable-features=TranslateUI',
                        '--disable-ipc-flooding-protection'
                    ]
                )
                
                try:
                    # Process all slides in parallel
                    # Create semaphore to limit concurrent operations
                    semaphore = asyncio.Semaphore(5)
                    context = await browser.new_context(
                        viewport={'width': 1920, 'height': 1080}
                    )
                    
                    context = await browser.new_context(
                        viewport={'width': 1920, 'height': 1080}
                    )
                    
                    async def process_single_slide(slide_info: Dict) -> Dict:
                        """Process a single slide with controlled concurrency."""
                        async with semaphore:
                            slide_num = slide_info['number']
                            
                            try:
                                # Create a new page for this slide
                                page = await context.new_page()
                                
                                # Set exact viewport dimensions
                                await page.set_viewport_size({"width": 1920, "height": 1080})
                                await page.emulate_media(media='screen')
                                
                                # Force device pixel ratio to 1
                                await page.evaluate(r"""
                                    () => {
                                        Object.defineProperty(window, 'devicePixelRatio', {
                                            get: () => 1
                                        });
                                    }
                                """)
                                
                                try:
                                    # Extract visual elements
                                    visual_elements = await self.extract_visual_elements(page, slide_info['path'], temp_path)
                                    
                                    # Capture clean background
                                    background_path = await self.capture_clean_background(page, slide_info['path'], temp_path, visual_elements)
                                    
                                    # Extract text elements
                                    text_elements = await self.extract_text_elements(page, slide_info['path'])
                                    
                                    slide_analysis = {
                                        'slide_info': slide_info,
                                        'visual_elements': visual_elements,
                                        'background_path': background_path,
                                        'text_elements': text_elements
                                    }
                                    
                                    return slide_analysis
                                    
                                except Exception as e:
                                    return {
                                        'slide_info': slide_info,
                                        'visual_elements': [],
                                        'background_path': None,
                                        'text_elements': [],
                                        'error': str(e)
                                    }
                                    
                                finally:
                                    # Always close the page to free memory
                                    await page.close()
                                    
                            except Exception as e:
                                return {
                                    'slide_info': slide_info,
                                    'visual_elements': [],
                                    'background_path': None,
                                    'text_elements': [],
                                    'error': f"Page creation failed: {str(e)}"
                                }
                    
                    # Launch ALL slides in parallel
                    parallel_tasks = [
                        process_single_slide(slide_info) 
                        for slide_info in self.slides_info
                    ]
                    
                    # Wait for ALL slides to complete in parallel
                    slide_analyses = await asyncio.gather(*parallel_tasks, return_exceptions=True)
                    
                    # Handle any top-level exceptions
                    processed_analyses = []
                    for i, result in enumerate(slide_analyses):
                        if isinstance(result, Exception):
                            error_analysis = {
                                'slide_info': self.slides_info[i],
                                'visual_elements': [],
                                'background_path': None,
                                'text_elements': [],
                                'error': str(result)
                            }
                            processed_analyses.append(error_analysis)
                        else:
                            processed_analyses.append(result)
                    
                    all_slide_analyses = processed_analyses
                    
                finally:
                    await browser.close()
            
            # Build PPTX presentation
            # Create new PowerPoint presentation
            presentation = Presentation()
            
            # Set slide dimensions to 1920x1080 (16:9)
            presentation.slide_width = Inches(20)  # 1920px at 96 DPI
            presentation.slide_height = Inches(11.25)  # 1080px at 96 DPI
            
            # Remove default slide
            if len(presentation.slides) > 0:
                xml_slides = presentation.slides._sldIdLst
                xml_slides.remove(xml_slides[0])
            
            # Build slides using the analyzed data
            successful_slides = 0
            for i, slide_analysis in enumerate(all_slide_analyses, 1):
                try:
                    if 'error' in slide_analysis and slide_analysis['error']:
                        # Create a blank slide with error message
                        blank_slide_layout = presentation.slide_layouts[6]
                        slide = presentation.slides.add_slide(blank_slide_layout)
                        
                        # Add error text box
                        textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(18), Inches(2))
                        text_frame = textbox.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = f"Error processing slide: {slide_analysis['error']}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(255, 0, 0)
                    else:
                        await self.build_slide_from_analysis(presentation, slide_analysis, temp_path)
                        successful_slides += 1
                        
                except Exception as e:
                    # Create error slide
                    blank_slide_layout = presentation.slide_layouts[6]
                    slide = presentation.slides.add_slide(blank_slide_layout)
                    
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(18), Inches(2))
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = f"Error building slide: {str(e)}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(255, 0, 0)
            
            # Save PowerPoint presentation
            presentation_name = self.metadata.get('presentation_name', 'presentation')
            temp_output_path = temp_path / f"{presentation_name}.pptx"
            
            presentation.save(str(temp_output_path))
            
            if store_locally:
                # Store in the static files directory for URL serving
                timestamp = int(asyncio.get_event_loop().time())
                filename = f"{presentation_name}_{timestamp}.pptx"
                final_output = output_dir / filename
                final_output.parent.mkdir(exist_ok=True)
                
                # Copy from temp to final location
                shutil.copy2(temp_output_path, final_output)
                
                return final_output, len(presentation.slides)
            else:
                # For direct download, read content from temp file before cleanup
                try:
                    with open(temp_output_path, 'rb') as f:
                        pptx_content = f.read()
                except Exception as e:
                    raise
                
                return pptx_content, len(presentation.slides), presentation_name


@router.post("/convert-to-pptx")
async def convert_presentation_to_pptx(request: ConvertRequest):
    """
    Convert HTML presentation to PPTX.
    
    Takes a presentation folder path and returns either:
    - PPTX file directly (if download=true) - uses presentation name as filename
    - JSON response with download URL (if download=false, default)
    """
    try:
        # Validate presentation path exists
        presentation_path = Path(request.presentation_path)
        if not presentation_path.exists():
            raise HTTPException(status_code=404, detail=f"Presentation path not found: {request.presentation_path}")
        
        # Check if metadata.json exists
        metadata_path = presentation_path / "metadata.json"
        if not metadata_path.exists():
            raise HTTPException(status_code=400, detail=f"metadata.json not found in: {request.presentation_path}")
        
        # Create converter
        converter = OptimizedHTMLToPPTXConverter(request.presentation_path)
        
        # If download is requested, don't store locally and return file directly
        if request.download:
            pptx_content, total_slides, presentation_name = await converter.convert_to_pptx(store_locally=False)
            
            return Response(
                content=pptx_content,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename=\"{presentation_name}.pptx\""}
            )
        
        # Otherwise, store locally and return JSON with download URL
        pptx_path, total_slides = await converter.convert_to_pptx(store_locally=True)
        
        pptx_url = f"/downloads/{pptx_path.name}"
        
        return ConvertResponse(
            success=True,
            message=f"PPTX generated successfully with {total_slides} slides",
            pptx_url=pptx_url,
            filename=pptx_path.name,
            total_slides=total_slides
        )
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX conversion failed: {str(e)}")


@router.get("/health")
async def pptx_health_check():
    """PPTX service health check endpoint."""
    return {
        "status": "healthy", 
        "service": "HTML to PPTX Converter"
    }